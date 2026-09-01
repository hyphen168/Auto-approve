"""PPT 演示文稿处理：AI 生成、已有 PPT 摘要。"""
import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import llm

PPT_SYSTEM = "你是专业的 PPT 大纲与文案设计专家，擅长为演示文稿设计清晰、有条理、有说服力的内容。"

ACCENT = RGBColor(0x1F, 0x38, 0x64)   # 深蓝主色
GREY = RGBColor(0x40, 0x40, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

RECT = 1  # MSO_SHAPE.RECTANGLE


def _shape_text(shape, text, size, bold=False, color=ACCENT, center=False):
    shape.text = text
    for p in shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def _add_body_tf(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def build_ppt_bytes(content, default_title="演示文稿"):
    """根据 AI 大纲内容构建 pptx，返回字节。

    解析规则：
    - 第一行 `# 大标题` 与第二行 `- 副标题` 组成封面
    - 之后每个 `## 页面标题` 开启一页，其后的 `- 要点` 为该页正文
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # Blank

    # ---- 解析大纲 ----
    cover = None
    cover_sub = None
    slides = []      # [(title, [bullets]), ...]
    cur_title = None
    cur_bullets = []

    for line in (content or "").splitlines():
        s = line.strip()
        if not s:
            continue
        # 标题（# / ## / ### ...）
        if s.startswith("#"):
            t = s.lstrip("#").strip()
            if cover is None:
                cover = t or None
            else:
                if cur_title is not None or cur_bullets:
                    slides.append((cur_title, cur_bullets))
                cur_title = t
                cur_bullets = []
            continue
        # 封面副标题：第一页 `# 大标题` 之后、第一个页面标题之前的 `- 副标题`
        if cover is not None and cur_title is None and cover_sub is None and s.startswith("-"):
            cover_sub = s.lstrip("-").strip() or None
            continue
        # 页面要点 / 普通行
        if s.startswith("-"):
            if cur_title is not None:
                cur_bullets.append(s.lstrip("-").strip())
            elif cover_sub is None:
                cover_sub = s.lstrip("-").strip() or None
        elif cur_title is not None:
            cur_bullets.append(s)
        elif cover_sub is None:
            cover_sub = s
    if cur_title is not None or cur_bullets:
        slides.append((cur_title, cur_bullets))

    if not cover:
        cover = default_title or "演示文稿"
    if not slides:
        slides.append((default_title or "内容", ["请在此填写要点"]))

    # ---- 封面页 ----
    s = prs.slides.add_slide(layout)
    bar = s.shapes.add_shape(RECT, Inches(0), Inches(0), prs.slide_width, Inches(2.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2))
    _shape_text(tb, cover, 40, True, WHITE, center=True)
    if cover_sub:
        tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.8))
        _shape_text(tb2, cover_sub, 22, False, WHITE, center=True)
    ft = _add_body_tf(s, 11.9, 6.95, 1.1, 0.4)
    p = ft.paragraphs[0]
    r = p.add_run()
    r.text = "1"
    r.font.size = Pt(12)
    r.font.color.rgb = GREY
    p.alignment = PP_ALIGN.RIGHT

    # ---- 正文页 ----
    for i, (title, bullets) in enumerate(slides, start=2):
        sl = prs.slides.add_slide(layout)
        bar = sl.shapes.add_shape(RECT, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.8))
        _shape_text(tb, title or "无标题", 28, True, WHITE)

        body = _add_body_tf(sl, 0.7, 1.5, 12, 5.4)
        body.vertical_anchor = MSO_ANCHOR.TOP
        for j, b in enumerate(bullets):
            p = body.paragraphs[0] if j == 0 else body.add_paragraph()
            r = p.add_run()
            r.text = "•  " + b
            r.font.size = Pt(20)
            r.font.color.rgb = GREY
            if len(b) <= 18:
                r.font.bold = True

        ft = _add_body_tf(sl, 11.9, 6.95, 1.1, 0.4)
        p = ft.paragraphs[0]
        r = p.add_run()
        r.text = str(i)
        r.font.size = Pt(12)
        r.font.color.rgb = GREY
        p.alignment = PP_ALIGN.RIGHT

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_ppt(cfg, topic, slide_count, extra=""):
    """根据主题生成 PPT，返回 (pptx_bytes, 原始AI大纲)。"""
    slide_count = max(2, int(slide_count))
    user = (
        "请为主题「%s」设计一份 %d 页的演示文稿大纲。\n"
        "附加要求：%s\n"
        "【输出格式要求】(严格遵守)\n"
        "1. 第一行：`# 大标题`（封面主标题）；\n"
        "2. 第二行：`- 副标题`（封面副标题，一句话概括）；\n"
        "3. 之后每一页：先写 `## 页面标题`，紧接着写至少 3 行要点，每行以 `- ` 开头；\n"
        "4. 最后一页为「总结」页；\n"
        "5. 直接输出大纲，不要输出任何解释。"
        % (topic, slide_count, extra or "无")
    )
    content = llm.chat(
        cfg,
        [{"role": "system", "content": PPT_SYSTEM},
         {"role": "user", "content": user}],
        max_tokens=4096,
        timeout=900,
    )
    return build_ppt_bytes(content, topic), content


def summarize_presentation(cfg, content):
    """对已有 PPT 生成内容摘要。"""
    prs = Presentation(io.BytesIO(content))
    shorts = []
    for slide in prs.slides:
        texts = [sh.text.strip() for sh in slide.shapes
                 if sh.has_text_frame and sh.text.strip()]
        shorts.append(" / ".join(texts))
    text = "\n".join(shorts)
    if len(text) > 8000:
        text = text[:8000] + "\n……（内容过长已截断）"
    user = (
        "请总结这份 PPT 的结构和内容，输出：\n"
        "1) 主题与目标；2) 各页核心要点；3) 改进建议。\n"
        "用「-」列表输出。\n\n【PPT 内容】\n" + text
    )
    return llm.chat(
        cfg,
        [{"role": "system", "content": PPT_SYSTEM},
         {"role": "user", "content": user}],
        max_tokens=2048,
        timeout=600,
    )